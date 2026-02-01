"""PowerPoint presentation processing and text extraction"""
from pptx import Presentation
from typing import Dict, List
import os


class PPTXProcessor:
    """Handles PowerPoint presentation processing and text extraction"""
    
    def __init__(self):
        self.supported_formats = ['.pptx', '.ppt']
    
    def extract_text(self, file_path: str) -> Dict[str, any]:
        """
        Extract text from PPTX file
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        try:
            prs = Presentation(file_path)
            
            # Extract metadata
            metadata = {
                'file_name': os.path.basename(file_path),
                'file_size': os.path.getsize(file_path),
                'num_slides': len(prs.slides)
            }
            
            # Extract text from all slides
            slides_content = []
            full_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_title = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text.append(text)
                        
                        # First text box is usually the title
                        if not slide_title and shape.has_text_frame:
                            slide_title = text
                
                slides_content.append({
                    'slide_number': slide_num,
                    'title': slide_title,
                    'content': slide_text,
                    'full_text': '\n'.join(slide_text)
                })
                
                full_text.append(f"Slide {slide_num}: {slide_title}\n" + '\n'.join(slide_text))
            
            return {
                'success': True,
                'full_text': '\n\n'.join(full_text),
                'slides': slides_content,
                'metadata': metadata,
                'format': 'pptx'
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'format': 'pptx'
            }
    
    def extract_with_structure(self, file_path: str) -> Dict[str, any]:
        """
        Extract text with structural information organized by slides
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Dictionary with structured content
        """
        result = self.extract_text(file_path)
        
        if not result['success']:
            return result
        
        # Structure content by slides as sections
        structured_content = []
        
        for slide in result['slides']:
            structured_content.append({
                'heading': slide['title'] or f"Slide {slide['slide_number']}",
                'level': f"Slide {slide['slide_number']}",
                'content': slide['content']
            })
        
        result['structured_content'] = structured_content
        return result
